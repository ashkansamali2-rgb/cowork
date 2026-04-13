"""
tools.py — Tool registry for the autonomous agent runtime.
Each tool is a callable that returns a str result.
"""
import ast
import asyncio
import json
import os
import subprocess
import sys
import time
from pathlib import Path
from typing import Any

AGENTS_MEM_DIR = Path.home() / "cowork" / "agents" / "memory"

# ── In-memory result cache ─────────────────────────────────────────────────────
_CACHE: dict = {}       # key -> (timestamp, result)
_CACHE_TTL: int = 600   # 10 minutes


def _cache_get(key: str):
    """Return cached value for key if still fresh, else None."""
    if key in _CACHE:
        ts, val = _CACHE[key]
        if time.time() - ts < _CACHE_TTL:
            return val
        del _CACHE[key]
    return None


def _cache_set(key: str, val: str) -> None:
    """Store val in cache under key with current timestamp."""
    _CACHE[key] = (time.time(), val)

# ── web_search(query: str) ─────────────────────────────────────────────────────
def web_search(query: str) -> str:
    """DuckDuckGo search. Returns top 5 results with titles, URLs, snippets."""
    if not query:
        return "Error: query is required"
    cached = _cache_get(query)
    if cached is not None:
        return cached
    try:
        try:
            from ddgs import DDGS  # new package name
        except ImportError:
            from duckduckgo_search import DDGS  # legacy fallback
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=5))
        if not results:
            return f"No results found for: {query}"
        lines = [f"Search results for: {query}\n"]
        for i, r in enumerate(results, 1):
            lines.append(f"{i}. {r.get('title', '')}")
            lines.append(f"   URL: {r.get('href', '')}")
            lines.append(f"   {r.get('body', '')[:300]}")
            lines.append("")
        result = "\n".join(lines)
        _cache_set(query, result)
        return result
    except Exception as e:
        return f"web_search failed: {e}"

# ── fetch_url(url: str) ────────────────────────────────────────────────────────
def fetch_url(url: str) -> str:
    """Fetch webpage and extract clean text using BeautifulSoup."""
    if not url:
        return "Error: url is required"
    cached = _cache_get(url)
    if cached is not None:
        return cached
    try:
        import requests
        from bs4 import BeautifulSoup
        headers = {"User-Agent": "Mozilla/5.0 (compatible; JarvisAgent/1.0)"}
        resp = requests.get(url, headers=headers, timeout=10)
        resp.raise_for_status()
        soup = BeautifulSoup(resp.text, "lxml")
        for tag in soup(["script", "style", "nav", "header", "footer", "aside"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        # Collapse blank lines
        import re
        text = re.sub(r'\n{3,}', '\n\n', text)
        result = text[:8000]
        _cache_set(url, result)
        return result
    except Exception as e:
        return f"fetch_url failed for {url}: {e}"

# ── run_shell(cmd: str) ────────────────────────────────────────────────────────
_FORBIDDEN = ["rm -rf /", "mkfs", "dd if=", ":(){:|:&};:", "sudo rm -rf /"]

def run_shell(cmd: str) -> str:
    """Execute a shell command and return stdout+stderr. Safety-checked."""
    for bad in _FORBIDDEN:
        if bad in cmd:
            return f"BLOCKED: dangerous command pattern detected: {bad}"
    try:
        result = subprocess.run(
            cmd, shell=True, capture_output=True, text=True, timeout=30
        )
        out = (result.stdout + result.stderr).strip()
        return out if out else "(command ran, no output)"
    except subprocess.TimeoutExpired:
        return "run_shell: command timed out after 30s"
    except Exception as e:
        return f"run_shell failed: {e}"

# ── read_file(path: str) ───────────────────────────────────────────────────────
def read_file(path: str) -> str:
    """Read file contents."""
    try:
        p = Path(os.path.expanduser(path))
        return p.read_text(errors="replace")[:10000]
    except Exception as e:
        return f"read_file failed: {e}"

# ── write_file(path: str, content: str) ───────────────────────────────────────
def write_file(path: str, content: str) -> str:
    """Write content to file, creating parent dirs as needed."""
    try:
        p = Path(os.path.expanduser(path))
        p.parent.mkdir(parents=True, exist_ok=True)
        p.write_text(content)
        return f"Written {len(content)} chars to {p}"
    except Exception as e:
        return f"write_file failed: {e}"

# ── append_file(path: str, content: str) ──────────────────────────────────────
def append_file(path: str, content: str) -> str:
    """Append content to file."""
    try:
        p = Path(os.path.expanduser(path))
        p.parent.mkdir(parents=True, exist_ok=True)
        with open(p, "a") as f:
            f.write(content)
        return f"Appended {len(content)} chars to {p}"
    except Exception as e:
        return f"append_file failed: {e}"

# ── list_dir(path: str) ────────────────────────────────────────────────────────
def list_dir(path: str = ".") -> str:
    """List directory contents."""
    try:
        p = Path(os.path.expanduser(path))
        entries = sorted(p.iterdir(), key=lambda x: (x.is_file(), x.name))
        lines = [f"{'[D]' if e.is_dir() else '[F]'} {e.name}" for e in entries]
        return f"Contents of {p}:\n" + "\n".join(lines) if lines else f"{p} is empty"
    except Exception as e:
        return f"list_dir failed: {e}"

# ── open_app(name: str) ────────────────────────────────────────────────────────
def open_app(name: str) -> str:
    """Open a macOS application and verify it opened."""
    try:
        import time as _time
        subprocess.run(["open", "-a", name], capture_output=True, text=True)
        _time.sleep(1)
        verify = subprocess.run(["pgrep", "-x", name], capture_output=True, text=True)
        if verify.returncode == 0:
            return f"Opened {name} successfully (pid: {verify.stdout.strip()})"
        return f"Opened {name} (could not verify via pgrep)"
    except Exception as e:
        return f"open_app failed: {e}"


# ── click_menu(app, menu, item) ────────────────────────────────────────────────
def click_menu(app: str, menu: str, item: str) -> str:
    """Click a menu item in a Mac application via osascript."""
    try:
        script = f'tell application "System Events" to tell process "{app}" to click menu item "{item}" of menu "{menu}" of menu bar 1'
        result = subprocess.run(["osascript", "-e", script], capture_output=True, text=True)
        if result.returncode == 0:
            return f"Clicked {app} > {menu} > {item}"
        return f"click_menu error: {result.stderr.strip()}"
    except Exception as e:
        return f"click_menu failed: {e}"


# ── type_text(text: str) ───────────────────────────────────────────────────────
def type_text(text: str) -> str:
    """Type text using keyboard via System Events."""
    try:
        escaped = text.replace('"', '\\"')
        script = f'tell application "System Events" to keystroke "{escaped}"'
        subprocess.run(["osascript", "-e", script])
        return f"Typed: {text[:50]}"
    except Exception as e:
        return f"type_text failed: {e}"


# ── press_key(key: str) ────────────────────────────────────────────────────────
def press_key(key: str) -> str:
    """Press a key or combo. Supports: return, tab, space, escape, delete, cmd+n, cmd+s, etc."""
    try:
        key_map = {"return": 36, "tab": 48, "space": 49, "escape": 53, "delete": 51}
        modifier_map = {"cmd": "command down", "shift": "shift down", "alt": "option down", "ctrl": "control down"}
        parts = key.lower().split("+")
        if len(parts) > 1:
            mods = ", ".join(modifier_map.get(p, f"{p} down") for p in parts[:-1])
            char = parts[-1]
            script = f'tell application "System Events" to keystroke "{char}" using {{{mods}}}'
        elif key in key_map:
            script = f'tell application "System Events" to key code {key_map[key]}'
        else:
            script = f'tell application "System Events" to keystroke "{key}"'
        subprocess.run(["osascript", "-e", script])
        return f"Pressed: {key}"
    except Exception as e:
        return f"press_key failed: {e}"


# ── get_screen_text() ─────────────────────────────────────────────────────────
def get_screen_text() -> str:
    """Take a screenshot and return the file path for vision analysis."""
    try:
        import time as _time
        path = f"/tmp/screenshot_{int(_time.time())}.png"
        subprocess.run(["screencapture", "-x", path])
        return f"Screenshot saved to {path}"
    except Exception as e:
        return f"get_screen_text failed: {e}"


# ── focus_app(name: str) ──────────────────────────────────────────────────────
def focus_app(name: str) -> str:
    """Bring a Mac application to the foreground."""
    try:
        subprocess.run(["osascript", "-e", f'tell application "{name}" to activate'])
        return f"Focused {name}"
    except Exception as e:
        return f"focus_app failed: {e}"


# ── create_keynote_presentation(title, slides) ────────────────────────────────
def create_keynote_presentation(title=None, slides=None, content=None) -> str:
    """Create a dark-themed PowerPoint presentation using python-pptx.
    Can be called as:
      create_keynote_presentation(title, slides)  where slides = list of {title, content} dicts
      create_keynote_presentation(content=string)  where string is parsed to extract title/slides
      create_keynote_presentation(string)          single positional string argument
    Falls back to osascript Keynote if python-pptx fails.
    """
    import re as _re

    # ── Normalise arguments ────────────────────────────────────────────────────
    # If called with a single positional string (title is actually the content)
    if title is not None and slides is None and content is None and isinstance(title, str) and '\n' in title:
        # Looks like full content was passed as title
        content = title
        title = None

    # Parse content string into title + slides list
    if content is not None and isinstance(content, str):
        lines = content.strip().splitlines()
        # Try to extract a title from the first non-empty line
        parsed_title = None
        for line in lines:
            stripped = line.strip()
            if stripped:
                parsed_title = _re.sub(r'^#+\s*', '', stripped)
                break
        title = parsed_title or "Presentation"
        # Build slides from sections (## headings or numbered items)
        slides = []
        current = None
        for line in lines[1:]:
            stripped = line.strip()
            if not stripped:
                continue
            if _re.match(r'^#{1,3}\s+', stripped) or _re.match(r'^\d+[\.\)]\s+', stripped):
                if current:
                    slides.append(current)
                slide_title = _re.sub(r'^#{1,3}\s+', '', stripped)
                slide_title = _re.sub(r'^\d+[\.\)]\s+', '', slide_title)
                current = {"title": slide_title, "content": ""}
            else:
                if current is None:
                    current = {"title": title, "content": ""}
                current["content"] = (current["content"] + "\n" + stripped).strip()
        if current:
            slides.append(current)
        if not slides:
            slides = [{"title": title, "content": content[:500]}]

    # Final defaults
    if title is None:
        title = "Presentation"
    if not slides:
        slides = [{"title": title, "content": ""}]

    # ── python-pptx primary method ─────────────────────────────────────────────
    try:
        from pptx import Presentation as _Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        COLOR_BG      = RGBColor(0x1a, 0x1a, 0x2e)
        COLOR_ACCENT  = RGBColor(0x7C, 0x3A, 0xED)
        COLOR_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
        COLOR_LIGHT   = RGBColor(0xCC, 0xCC, 0xFF)

        prs = _Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]  # completely blank

        def _set_bg(slide):
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = COLOR_BG

        def _add_rect(slide, left, top, width, height, color):
            from pptx.util import Emu
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()  # no border

        def _add_textbox(slide, text, left, top, width, height, font_size=24,
                         bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
            txb = slide.shapes.add_textbox(left, top, width, height)
            tf  = txb.text_frame
            tf.word_wrap = wrap
            p   = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color or COLOR_WHITE

        # Title slide
        title_slide = prs.slides.add_slide(blank_layout)
        _set_bg(title_slide)
        # Accent bar
        _add_rect(title_slide, Inches(0), Inches(3.2), Inches(13.33), Inches(0.08), COLOR_ACCENT)
        _add_textbox(title_slide, title,
                     Inches(1), Inches(2.2), Inches(11.33), Inches(1.5),
                     font_size=44, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
        _add_textbox(title_slide, f"{len(slides)} topics",
                     Inches(1), Inches(3.5), Inches(11.33), Inches(0.6),
                     font_size=20, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)

        # Content slides
        for slide_data in slides:
            sl = prs.slides.add_slide(blank_layout)
            _set_bg(sl)
            # Top accent bar
            _add_rect(sl, Inches(0), Inches(0), Inches(13.33), Inches(0.06), COLOR_ACCENT)
            slide_title   = slide_data.get("title", "")
            slide_content = slide_data.get("content", "")
            _add_textbox(sl, slide_title,
                         Inches(0.5), Inches(0.2), Inches(12.33), Inches(1.0),
                         font_size=32, bold=True, color=COLOR_ACCENT)
            if slide_content:
                # Split into bullet lines
                bullet_lines = [ln.strip() for ln in slide_content.splitlines() if ln.strip()]
                display_text = "\n".join(f"• {ln}" for ln in bullet_lines[:12])
                _add_textbox(sl, display_text,
                             Inches(0.5), Inches(1.4), Inches(12.33), Inches(5.5),
                             font_size=20, color=COLOR_LIGHT)

        safe_title = _re.sub(r'[^a-zA-Z0-9_\- ]', '', title).strip().replace(' ', '-') or 'Presentation'
        output_path = os.path.expanduser(f"~/Desktop/{safe_title}.pptx")
        prs.save(output_path)
        os.system(f'open "{output_path}"')
        return f"PowerPoint presentation created and opened: {output_path}"

    except Exception as pptx_err:
        # ── osascript Keynote fallback ─────────────────────────────────────────
        try:
            import os as _os
            output_path = _os.path.expanduser(f"~/Desktop/{title}.key")
            slide_lines = []
            for i, slide in enumerate(slides):
                st = slide.get('title', f'Slide {i+1}').replace('"', '\\"')
                sc = slide.get('content', '').replace('"', '\\"').replace('\n', '\\n')
                if i == 0:
                    slide_lines.append(f'set the object text of the default title item of slide 1 to "{st}"')
                    if sc:
                        slide_lines.append(f'set the object text of the default body item of slide 1 to "{sc}"')
                else:
                    slide_lines.append(f'set ns to duplicate slide 1 to end of slides')
                    slide_lines.append(f'set the object text of the default title item of ns to "{st}"')
                    if sc:
                        slide_lines.append(f'set the object text of the default body item of ns to "{sc}"')
            slides_script = '\n'.join(slide_lines)
            script = f'''tell application "Keynote"
    activate
    set newDoc to make new document with properties {{document theme:theme "White"}}
    tell newDoc
        {slides_script}
    end tell
    save newDoc in POSIX file "{output_path}"
end tell'''
            result = subprocess.run(["osascript", "-e", script], capture_output=True, text=True)
            if result.returncode == 0:
                return f"Keynote presentation created (pptx failed: {pptx_err}): {output_path}"
            return f"create_keynote_presentation pptx error: {pptx_err}; osascript error: {result.stderr.strip()}"
        except Exception as e:
            return f"create_keynote_presentation failed entirely: pptx={pptx_err}, keynote={e}"


# ── create_pages_document(title, content) ─────────────────────────────────────
def create_pages_document(title: str, content: str) -> str:
    """Create a Pages document on the Desktop."""
    try:
        import os as _os
        output_path = _os.path.expanduser(f"~/Desktop/{title}.pages")
        safe_content = content.replace('"', '\\"').replace('\n', '\\n')
        script = f'''tell application "Pages"
    activate
    set newDoc to make new document
    tell body text of newDoc
        set its paragraphs to "{safe_content}"
    end tell
    save newDoc in POSIX file "{output_path}"
end tell'''
        result = subprocess.run(["osascript", "-e", script], capture_output=True, text=True)
        if result.returncode == 0:
            return f"Pages document created: {output_path}"
        return f"create_pages_document error: {result.stderr.strip()}"
    except Exception as e:
        return f"create_pages_document failed: {e}"


# ── send_notification(title, message) ────────────────────────────────────────
def send_notification(title: str, message: str) -> str:
    """Send a macOS notification banner."""
    try:
        safe_msg = message.replace('"', '\\"')
        safe_title = title.replace('"', '\\"')
        subprocess.run(["osascript", "-e", f'display notification "{safe_msg}" with title "{safe_title}"'])
        return f"Notification sent: {title}"
    except Exception as e:
        return f"send_notification failed: {e}"


# ── clipboard_paste() ─────────────────────────────────────────────────────────
def clipboard_paste() -> str:
    """Paste clipboard contents at current cursor position."""
    try:
        subprocess.run(["osascript", "-e", 'tell application "System Events" to keystroke "v" using command down'])
        return "Pasted from clipboard"
    except Exception as e:
        return f"clipboard_paste failed: {e}"


# ── move_mouse(x, y) ─────────────────────────────────────────────────────────
def move_mouse(x: int, y: int) -> str:
    """Move mouse cursor to screen coordinates."""
    try:
        import pyautogui
        pyautogui.moveTo(x, y)
        return f"Moved mouse to ({x}, {y})"
    except ImportError:
        subprocess.run([sys.executable, "-m", "pip", "install", "pyautogui", "pillow", "-q"])
        return "pyautogui not installed — installing now, retry next step"
    except Exception as e:
        return f"move_mouse failed: {e}"


# ── click_at(x, y) ───────────────────────────────────────────────────────────
def click_at(x: int, y: int) -> str:
    """Click at screen coordinates."""
    try:
        import pyautogui
        pyautogui.click(x, y)
        return f"Clicked at ({x}, {y})"
    except ImportError:
        return "pyautogui not available — run move_mouse first to install"
    except Exception as e:
        return f"click_at failed: {e}"


# ── screenshot_region(x, y, w, h) ────────────────────────────────────────────
def screenshot_region(x: int, y: int, w: int, h: int) -> str:
    """Take a screenshot of a specific screen region."""
    try:
        import time as _time
        path = f"/tmp/region_{int(_time.time())}.png"
        subprocess.run(["screencapture", "-x", "-R", f"{x},{y},{w},{h}", path])
        return f"Region screenshot saved to {path}"
    except Exception as e:
        return f"screenshot_region failed: {e}"

# ── speak(text: str) ──────────────────────────────────────────────────────────
def speak(text: str) -> str:
    """Speak text via Jarvis TTS (bus TASK_VOICE event)."""
    try:
        import asyncio as _asyncio
        async def _send():
            import websockets as _ws
            async with _ws.connect("ws://127.0.0.1:8002") as ws:
                await ws.send(json.dumps({"register": "agent-tools"}))
                await ws.recv()
                await ws.send(json.dumps({"type": "TASK_VOICE", "msg": text[:500]}))
        _asyncio.run(_send())
        return f"Spoke: {text[:100]}"
    except Exception as e:
        return f"speak failed (TTS may be offline): {e}"

# ── create_document(title, content, fmt) ─────────────────────────────────────
def create_document(title: str, content: str, fmt: str = "md") -> str:
    """Create a document on the Desktop."""
    try:
        import re
        safe = re.sub(r'[^a-zA-Z0-9_\- ]', '', title).strip().replace(" ", "-")
        fname = f"{safe}.{fmt}"
        desktop = Path.home() / "Desktop" / fname
        desktop.write_text(content)
        return f"Document created: {desktop} ({len(content)} chars)"
    except Exception as e:
        return f"create_document failed: {e}"

# ── create_word_document(title, content, path) ────────────────────────────────
def create_word_document(title: str, content: str, path: str = None) -> str:
    """Create a properly formatted Word (.docx) document on the Desktop.
    Converts markdown headings (##) to Word heading styles and bullet points
    to proper Word list items. Opens the file automatically."""
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        import re
        from datetime import datetime

        doc = Document()

        # Set document title style
        title_para = doc.add_heading(title, level=0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.LEFT

        # Parse and add content line by line
        for line in content.splitlines():
            stripped = line.rstrip()

            # h3 ###
            if stripped.startswith("### "):
                doc.add_heading(stripped[4:], level=3)

            # h2 ##
            elif stripped.startswith("## "):
                doc.add_heading(stripped[3:], level=2)

            # h1 #
            elif stripped.startswith("# "):
                doc.add_heading(stripped[2:], level=1)

            # Bullet points: - or * or •
            elif re.match(r'^[-*•]\s+', stripped):
                text = re.sub(r'^[-*•]\s+', '', stripped)
                p = doc.add_paragraph(style='List Bullet')
                p.add_run(text)

            # Numbered list: 1. 2. etc
            elif re.match(r'^\d+\.\s+', stripped):
                text = re.sub(r'^\d+\.\s+', '', stripped)
                p = doc.add_paragraph(style='List Number')
                p.add_run(text)

            # Horizontal rule
            elif stripped in ('---', '***', '___'):
                doc.add_paragraph('─' * 60)

            # Empty line
            elif not stripped:
                doc.add_paragraph('')

            # Regular text
            else:
                # Handle inline bold **text**
                p = doc.add_paragraph()
                parts = re.split(r'\*\*(.*?)\*\*', stripped)
                for i, part in enumerate(parts):
                    run = p.add_run(part)
                    if i % 2 == 1:  # odd parts are bold
                        run.bold = True

        # Determine output path
        if path:
            out = Path(os.path.expanduser(path))
        else:
            import re as _re
            safe = _re.sub(r'[^a-zA-Z0-9_\- ]', '', title).strip().replace(" ", "-")
            date_str = datetime.now().strftime("%Y%m%d")
            out = Path.home() / "Desktop" / f"{safe}-{date_str}.docx"

        out.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(out))
        os.system(f"open \"{out}\"")
        return f"Word document created and opened: {out}"
    except Exception as e:
        return f"create_word_document failed: {e}"

# ── take_screenshot() ─────────────────────────────────────────────────────────
def take_screenshot() -> str:
    """Take a screenshot and save to Desktop."""
    from datetime import datetime
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    path = Path.home() / "Desktop" / f"screenshot_{ts}.png"
    result = subprocess.run(["screencapture", "-x", str(path)], capture_output=True)
    return str(path) if result.returncode == 0 else "screenshot failed"

# ── get_clipboard() ───────────────────────────────────────────────────────────
def get_clipboard() -> str:
    """Get clipboard contents."""
    try:
        result = subprocess.run(["pbpaste"], capture_output=True, text=True)
        return result.stdout or "(clipboard empty)"
    except Exception as e:
        return f"get_clipboard failed: {e}"

# ── set_clipboard(text: str) ──────────────────────────────────────────────────
def set_clipboard(text: str) -> str:
    """Set clipboard contents."""
    try:
        proc = subprocess.Popen(["pbcopy"], stdin=subprocess.PIPE)
        proc.communicate(text.encode())
        return "Clipboard set."
    except Exception as e:
        return f"set_clipboard failed: {e}"

# ── http_request(url, method, data) ───────────────────────────────────────────
def http_request(url: str, method: str = "GET", data: dict = None) -> str:
    """Make an HTTP request and return response text."""
    try:
        import requests
        method = method.upper()
        resp = requests.request(method, url, json=data, timeout=15)
        return resp.text[:5000]
    except Exception as e:
        return f"http_request failed: {e}"

# ── summarize(text, instruction) ─────────────────────────────────────────────
def summarize(text: str, instruction: str = "Summarize concisely.") -> str:
    """Use Qwen to summarize or transform text."""
    try:
        import requests
        payload = {
            "model": "qwen",
            "messages": [
                {"role": "system", "content": instruction},
                {"role": "user", "content": text[:6000]},
            ],
            "temperature": 0.3,
            "max_tokens": 512,
        }
        resp = requests.post("http://localhost:8081/v1/chat/completions", json=payload, timeout=60)
        resp.raise_for_status()
        return resp.json()["choices"][0]["message"]["content"].strip()
    except Exception as e:
        return f"summarize failed (Qwen may be offline): {e}"

# ── remember(key, value) ──────────────────────────────────────────────────────
def remember(key: str, value: str) -> str:
    """Store a key-value pair in persistent agent memory."""
    try:
        AGENTS_MEM_DIR.mkdir(parents=True, exist_ok=True)
        p = AGENTS_MEM_DIR / f"{key}.json"
        p.write_text(json.dumps({"key": key, "value": value}))
        return f"Remembered: {key}"
    except Exception as e:
        return f"remember failed: {e}"

# ── recall(key) ───────────────────────────────────────────────────────────────
def recall(key: str) -> str:
    """Retrieve a value from persistent agent memory."""
    try:
        p = AGENTS_MEM_DIR / f"{key}.json"
        if not p.exists():
            return f"No memory found for key: {key}"
        data = json.loads(p.read_text())
        return str(data.get("value", ""))
    except Exception as e:
        return f"recall failed: {e}"

# ── spawn_subagent(task) ──────────────────────────────────────────────────────
def spawn_subagent(task: str) -> str:
    """Spawn a child agent for a subtask. Returns agent_id."""
    try:
        import time
        agent_id = f"SUB-{int(time.time())}"
        import asyncio as _asyncio
        async def _send():
            import websockets as _ws
            async with _ws.connect("ws://127.0.0.1:8002") as ws:
                await ws.send(json.dumps({"register": "subagent-spawner"}))
                await ws.recv()
                await ws.send(json.dumps({
                    "type": "TASK_RESEARCH",
                    "msg": task,
                    "agent_id": agent_id,
                    "cwd": str(Path.home() / "cowork"),
                }))
        _asyncio.run(_send())
        return f"Spawned subagent {agent_id} for: {task}"
    except Exception as e:
        return f"spawn_subagent failed: {e}"

# ── read_screen(question) ─────────────────────────────────────────────────────
async def read_screen(question: str = "What do you see?") -> str:
    """Take a screenshot and ask Gemma 4 what is on screen."""
    try:
        import sys
        sys.path.insert(0, str(Path.home() / "cowork" / "jarvis"))
        from core.vision.screen_reader import ScreenReader
        return await ScreenReader().understand(question)
    except Exception as e:
        return f"[read_screen error: {e}]"


# ── find_and_click(description) ───────────────────────────────────────────────
async def find_and_click(description: str) -> str:
    """Find a UI element on screen and click it."""
    try:
        import sys
        sys.path.insert(0, str(Path.home() / "cowork" / "jarvis"))
        from core.vision.screen_reader import ScreenReader
        coords = await ScreenReader().find_element(description)
        if "NOT_FOUND" in coords or "x=" not in coords:
            return f"Element not found: {description}"
        import re
        m = re.search(r"x=(\d+).*?y=(\d+)", coords)
        if not m:
            return f"Could not parse coordinates: {coords}"
        x, y = int(m.group(1)), int(m.group(2))
        import pyautogui
        pyautogui.click(x, y)
        return f"Clicked at ({x}, {y}) for: {description}"
    except Exception as e:
        return f"[find_and_click error: {e}]"


# ── fix_screen_error() ────────────────────────────────────────────────────────
async def fix_screen_error() -> str:
    """Read any error on screen and attempt to diagnose it."""
    try:
        import sys
        sys.path.insert(0, str(Path.home() / "cowork" / "jarvis"))
        from core.vision.screen_reader import ScreenReader
        error_text = await ScreenReader().read_error_on_screen()
        if error_text == "NONE":
            return "No error visible on screen."
        return f"Error detected: {error_text}"
    except Exception as e:
        return f"[fix_screen_error: {e}]"


# ── Browser tools (via osascript / Safari) ────────────────────────────────────

def browser_navigate(url: str) -> str:
    """Open a URL in Safari."""
    import subprocess
    subprocess.run(["osascript", "-e",
        f'tell application "Safari" to open location "{url}"'], timeout=10)
    return f"Navigated to {url}"


def browser_get_text() -> str:
    """Get visible text from current Safari page."""
    import subprocess
    r = subprocess.run(["osascript", "-e", """
        tell application "Safari"
            set pageText to do JavaScript "document.body.innerText.substring(0, 5000)" in document 1
            return pageText
        end tell
    """], capture_output=True, text=True, timeout=10)
    return r.stdout.strip() or "(no text)"


def browser_click_link(link_text: str) -> str:
    """Click a link in Safari by its text."""
    import subprocess
    safe = link_text.replace("'", "\\'")
    subprocess.run(["osascript", "-e", f"""
        tell application "Safari"
            do JavaScript "var links=document.querySelectorAll('a');for(var l of links){{if(l.innerText.includes('{safe}')){{l.click();break;}}}}" in document 1
        end tell
    """], timeout=10)
    return f"Clicked link: {link_text}"


def browser_screenshot() -> str:
    """Screenshot current Safari page."""
    import subprocess, time
    path = f"/tmp/browser_{int(time.time())}.png"
    subprocess.run(["screencapture", "-x", path], timeout=5)
    return f"Screenshot saved: {path}"


# ── File operation tools ───────────────────────────────────────────────────────

import glob as _glob


def find_files(pattern: str, directory: str = ".") -> str:
    """Find files matching a glob pattern."""
    import os
    base = os.path.expanduser(directory)
    matches = _glob.glob(os.path.join(base, "**", pattern), recursive=True)
    matches = [m for m in matches if ".git" not in m][:20]
    return "\n".join(matches) if matches else "(no matches)"


def grep_content(pattern: str, directory: str = ".") -> str:
    """Search file contents for a pattern."""
    import subprocess, os
    base = os.path.expanduser(directory)
    r = subprocess.run(["grep", "-r", "--include=*.py", "--include=*.js",
                        "-l", pattern, base], capture_output=True, text=True, timeout=10)
    lines = [l for l in r.stdout.splitlines() if ".git" not in l][:10]
    return "\n".join(lines) if lines else "(not found)"


def zip_folder(path: str, output: str = None) -> str:
    """Zip a directory."""
    import subprocess, os, time
    path = os.path.expanduser(path)
    if not output:
        output = f"{path}_{int(time.time())}.zip"
    r = subprocess.run(["zip", "-r", output, path], capture_output=True, text=True, timeout=30)
    if r.returncode == 0:
        return f"Zipped to: {output}"
    return f"Zip failed: {r.stderr}"


# ── System control tools ───────────────────────────────────────────────────────

def get_running_apps() -> str:
    """List currently running apps."""
    import subprocess
    r = subprocess.run(["osascript", "-e",
        'tell application "System Events" to get name of every process whose background only is false'],
        capture_output=True, text=True, timeout=10)
    return r.stdout.strip()


def get_frontmost_app() -> str:
    """Get the currently active application name."""
    import subprocess
    r = subprocess.run(["osascript", "-e",
        'tell application "System Events" to get name of first process whose frontmost is true'],
        capture_output=True, text=True, timeout=5)
    return r.stdout.strip()


def close_app(app: str) -> str:
    """Quit an application."""
    import subprocess
    subprocess.run(["osascript", "-e", f'tell application "{app}" to quit'], timeout=10)
    return f"Closed: {app}"


# ── Code execution tools ───────────────────────────────────────────────────────

def run_python_code(code: str) -> str:
    """Run Python code in subprocess, return output. 30s timeout."""
    import subprocess, tempfile, os
    with tempfile.NamedTemporaryFile(mode="w", suffix=".py", delete=False) as f:
        f.write(code)
        tmp = f.name
    try:
        r = subprocess.run(["python3", tmp], capture_output=True, text=True, timeout=30)
        out = (r.stdout + r.stderr).strip()
        return out[:2000] if out else "(no output)"
    except subprocess.TimeoutExpired:
        return "[timeout after 30s]"
    finally:
        os.unlink(tmp)


def run_shell_safe(cmd: str, cwd: str = None) -> str:
    """Run shell command safely. 30s timeout."""
    import subprocess, os
    cwd = os.path.expanduser(cwd) if cwd else None
    try:
        r = subprocess.run(cmd, shell=True, capture_output=True, text=True,
                          timeout=30, cwd=cwd)
        out = (r.stdout + r.stderr).strip()
        return out[:2000] if out else "(no output)"
    except subprocess.TimeoutExpired:
        return "[timeout after 30s]"


# ── Tool registry ─────────────────────────────────────────────────────────────

TOOLS: dict[str, callable] = {
    "web_search":                   web_search,
    "fetch_url":                    fetch_url,
    "run_shell":                    run_shell,
    "read_file":                    read_file,
    "write_file":                   write_file,
    "append_file":                  append_file,
    "list_dir":                     list_dir,
    "open_app":                     open_app,
    "speak":                        speak,
    "create_document":              create_document,
    "create_word_document":         create_word_document,
    "take_screenshot":              take_screenshot,
    "get_clipboard":                get_clipboard,
    "set_clipboard":                set_clipboard,
    "http_request":                 http_request,
    "summarize":                    summarize,
    "remember":                     remember,
    "recall":                       recall,
    "spawn_subagent":               spawn_subagent,
    "click_menu":                   click_menu,
    "type_text":                    type_text,
    "press_key":                    press_key,
    "get_screen_text":              get_screen_text,
    "focus_app":                    focus_app,
    "create_keynote_presentation":  create_keynote_presentation,
    "create_pages_document":        create_pages_document,
    "send_notification":            send_notification,
    "clipboard_paste":              clipboard_paste,
    "move_mouse":                   move_mouse,
    "click_at":                     click_at,
    "screenshot_region":            screenshot_region,
    "read_screen":                  read_screen,
    "find_and_click":               find_and_click,
    "fix_screen_error":             fix_screen_error,
}

TOOL_DESCRIPTIONS: dict[str, str] = {
    "web_search":                   "web_search(query) — DuckDuckGo search, returns top 5 results with snippets",
    "fetch_url":                    "fetch_url(url) — Fetch webpage, extract clean text (up to 8000 chars)",
    "run_shell":                    "run_shell(cmd) — Run a shell command, returns stdout+stderr",
    "read_file":                    "read_file(path) — Read file contents",
    "write_file":                   "write_file(path, content) — Write content to file (creates dirs)",
    "append_file":                  "append_file(path, content) — Append content to existing file",
    "list_dir":                     "list_dir(path) — List directory contents",
    "open_app":                     "open_app(name) — Open a macOS application by name, verifies it opened",
    "speak":                        "speak(text) — Speak text aloud via Jarvis TTS",
    "create_document":              "create_document(title, content, fmt='md') — Create plain text/markdown document on Desktop",
    "create_word_document":         "create_word_document(title, content, path=None) — Create formatted Word .docx on Desktop, opens it automatically. PREFER THIS over create_document when user says 'create a document', 'write it up', 'document it', or 'write a report'.",
    "take_screenshot":              "take_screenshot() — Take screenshot, saves to Desktop, returns path",
    "get_clipboard":                "get_clipboard() — Get clipboard text contents",
    "set_clipboard":                "set_clipboard(text) — Copy text to clipboard",
    "http_request":                 "http_request(url, method='GET', data=None) — HTTP request, returns response",
    "summarize":                    "summarize(text, instruction) — Use Qwen LLM to summarize or transform text",
    "remember":                     "remember(key, value) — Persist key-value to ~/cowork/agents/memory/",
    "recall":                       "recall(key) — Retrieve persisted value by key",
    "spawn_subagent":               "spawn_subagent(task) — Spawn a child agent for a parallel subtask",
    "click_menu":                   "click_menu(app, menu, item) — Click a menu item in any Mac app via osascript",
    "type_text":                    "type_text(text) — Type text at current cursor position via keyboard",
    "press_key":                    "press_key(key) — Press a key or combo (return, tab, space, escape, cmd+n, cmd+s, cmd+v, etc.)",
    "get_screen_text":              "get_screen_text() — Take a screenshot and return the file path for vision analysis",
    "focus_app":                    "focus_app(name) — Bring a Mac application to the foreground",
    "create_keynote_presentation":  "create_keynote_presentation(title, slides) — Create a dark-themed PowerPoint .pptx on Desktop (opens automatically). slides = list of {title, content} dicts. Also accepts a single string content argument.",
    "create_pages_document":        "create_pages_document(title, content) — Create a Pages document on the Desktop",
    "send_notification":            "send_notification(title, message) — Show a macOS notification banner",
    "clipboard_paste":              "clipboard_paste() — Paste clipboard contents at current cursor (Cmd+V)",
    "move_mouse":                   "move_mouse(x, y) — Move mouse cursor to screen coordinates",
    "click_at":                     "click_at(x, y) — Click at screen coordinates",
    "screenshot_region":            "screenshot_region(x, y, w, h) — Take screenshot of a specific screen region, returns path",
    "read_screen":                  "read_screen(question='What do you see?') — Take a screenshot and ask Gemma 4 multimodal what is on screen",
    "find_and_click":               "find_and_click(description) — Find a UI element on screen by description and click it",
    "fix_screen_error":             "fix_screen_error() — Read any visible error on screen and return a diagnosis",
    "create_spreadsheet":           "create_spreadsheet(title, headers, rows, filename=None) — Create a styled xlsx spreadsheet and open it on Desktop",
    "create_pdf":                   "create_pdf(title, content, filename=None) — Create a formatted PDF document and open it on Desktop",
    "analyze_image":                "analyze_image(path, question='What do you see in this image?') — Send an image file to Gemma 4 multimodal for analysis",
    "analyze_screenshot":           "analyze_screenshot(question='What do you see?') — Take a screenshot and analyze it with Gemma 4 multimodal",
}

def get_tool_descriptions() -> str:
    return "\n".join(f"- {v}" for v in TOOL_DESCRIPTIONS.values())

def register_tool(name: str, fn: callable, description: str = ""):
    """Dynamically register a new tool at runtime."""
    TOOLS[name] = fn
    TOOL_DESCRIPTIONS[name] = description or f"{name}(...) — dynamically loaded skill"

# Meta-skill: create new skills on the fly
try:
    from core.agents.skills.skill_creator import skill_creator
    TOOLS["skill_creator"] = skill_creator
    TOOL_DESCRIPTIONS["skill_creator"] = "Create a brand new tool/skill when you need a capability that doesn't exist. Args: task_description (str)"
except ImportError:
    pass

# ── Document tools ────────────────────────────────────────────────────────────
def create_spreadsheet(title: str, headers: list, rows: list, filename: str = None) -> str:
    """Create an xlsx spreadsheet and open it."""
    try:
        import openpyxl
        from openpyxl.styles import PatternFill, Font, Alignment
        from datetime import date
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = title[:31]
        header_fill = PatternFill("solid", fgColor="7C3AED")
        header_font = Font(color="FFFFFF", bold=True)
        for col_idx, h in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col_idx, value=h)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal="center")
        alt_fill = PatternFill("solid", fgColor="F4F3F0")
        for row_idx, row in enumerate(rows, 2):
            fill = alt_fill if row_idx % 2 == 0 else None
            for col_idx, val in enumerate(row, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=val)
                if fill:
                    cell.fill = fill
        for col in ws.columns:
            max_len = max((len(str(c.value or "")) for c in col), default=10)
            ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 50)
        out = filename or os.path.expanduser(f"~/Desktop/{title.replace(' ','_')}_{date.today()}.xlsx")
        wb.save(out)
        subprocess.run(["open", out])
        return f"Spreadsheet saved: {out}"
    except ImportError:
        return "Error: openpyxl not installed. Run: pip install openpyxl"
    except Exception as e:
        return f"Spreadsheet error: {e}"


def create_pdf(title: str, content: str, filename: str = None) -> str:
    """Create a PDF document and open it."""
    try:
        from fpdf import FPDF
        from datetime import date
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 18)
        pdf.cell(0, 12, title, ln=True, align="C")
        pdf.ln(4)
        pdf.set_draw_color(124, 58, 237)
        pdf.set_line_width(0.5)
        pdf.line(10, pdf.get_y(), 200, pdf.get_y())
        pdf.ln(6)
        pdf.set_font("Helvetica", size=11)
        for line in content.split("\n"):
            stripped = line.strip()
            if not stripped:
                pdf.ln(4)
                continue
            if stripped.startswith("## "):
                pdf.set_font("Helvetica", "B", 13)
                pdf.cell(0, 8, stripped[3:], ln=True)
                pdf.set_font("Helvetica", size=11)
            elif stripped.startswith("# "):
                pdf.set_font("Helvetica", "B", 15)
                pdf.cell(0, 10, stripped[2:], ln=True)
                pdf.set_font("Helvetica", size=11)
            else:
                pdf.multi_cell(0, 6, stripped)
        pdf.set_font("Helvetica", "I", 9)
        pdf.set_y(-15)
        pdf.cell(0, 10, f"Page {pdf.page_no()}", align="C")
        out = filename or os.path.expanduser(f"~/Desktop/{title.replace(' ','_')}_{date.today()}.pdf")
        pdf.output(out)
        subprocess.run(["open", out])
        return f"PDF saved: {out}"
    except ImportError:
        return "Error: fpdf2 not installed. Run: pip install fpdf2"
    except Exception as e:
        return f"PDF error: {e}"


def analyze_image(path: str, question: str = "What do you see in this image?") -> str:
    """Send an image to Gemma 4 multimodal for analysis."""
    import base64
    import requests as _req
    try:
        with open(os.path.expanduser(path), "rb") as f:
            b64 = base64.b64encode(f.read()).decode()
        resp = _req.post("http://localhost:8080/v1/chat/completions", json={
            "messages": [{"role": "user", "content": [
                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}},
                {"type": "text", "text": question}
            ]}],
            "max_tokens": 500,
            "stream": False,
        }, timeout=30)
        return resp.json()["choices"][0]["message"]["content"]
    except Exception as e:
        return f"Image analysis error: {e}"


def analyze_screenshot(question: str = "What do you see?") -> str:
    """Take a screenshot and analyze it with Gemma."""
    import time as _time
    path = f"/tmp/screenshot_{int(_time.time())}.png"
    subprocess.run(["screencapture", "-x", path])
    return analyze_image(path, question)


TOOLS["create_spreadsheet"] = create_spreadsheet
TOOLS["create_pdf"] = create_pdf
TOOLS["analyze_image"] = analyze_image
TOOLS["analyze_screenshot"] = analyze_screenshot

# ── Browser automation tools ──────────────────────────────────────────────────
try:
    from core.agents.browser_agent import (
        browser_navigate, browser_get_page_text, browser_get_current_url,
        browser_click, browser_fill_form, browser_screenshot, browser_scroll_down,
    )
    TOOLS["browser_navigate"] = browser_navigate
    TOOLS["browser_get_page_text"] = browser_get_page_text
    TOOLS["browser_get_current_url"] = browser_get_current_url
    TOOLS["browser_click"] = browser_click
    TOOLS["browser_fill_form"] = browser_fill_form
    TOOLS["browser_screenshot"] = browser_screenshot
    TOOLS["browser_scroll_down"] = browser_scroll_down
except ImportError:
    pass

# ── Project builder tools ──────────────────────────────────────────────────────
try:
    from core.agents.project_builder import (
        build_project as _build_project_async,
        launch_project as _launch_project_async,
    )

    async def build_project(
        name: str,
        project_type: str,
        description: str = "",
        location: str = "~/Desktop",
        port: int = 3000,
    ) -> str:
        """Scaffold, implement, and optionally launch a project from a template."""
        return await _build_project_async(
            name=name,
            project_type=project_type,
            description=description,
            location=location,
            port=port,
        )

    async def launch_project(project_path: str, port: int = 3000) -> str:
        """Launch an existing project in Terminal and open it in the browser."""
        return await _launch_project_async(project_path=project_path, port=port)

    TOOLS["build_project"] = build_project
    TOOLS["launch_project"] = launch_project
    TOOL_DESCRIPTIONS["build_project"] = (
        "build_project(name, project_type, description='', location='~/Desktop', port=3000) — "
        "Scaffold a new project from a template (fastapi, flask, python, express), "
        "install dependencies, and optionally implement functionality via Qwen."
    )
    TOOL_DESCRIPTIONS["launch_project"] = (
        "launch_project(project_path, port=3000) — "
        "Launch an existing project in a Terminal window and open it in the browser."
    )
except ImportError:
    pass
